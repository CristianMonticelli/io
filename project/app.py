from flask import Flask, render_template, request, send_file
from pptx import Presentation
import io

app = Flask(__name__)

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    title = request.form.get("title")
    content = request.form.get("content")

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = content

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return send_file(
        file_stream,
        download_name="output.pptx",
        as_attachment=True
    )

if __name__ == "__main__":
    app.run(debug=True)